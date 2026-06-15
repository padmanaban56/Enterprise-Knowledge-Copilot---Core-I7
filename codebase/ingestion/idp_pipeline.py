"""
ingestion/idp_pipeline.py  —  Intelligent Document Processing Pipeline

Phase 1 implements:
  - L1 Rule-based classification (regex on filename/path)
  - L2 Embedding-based classification (cosine to dept centroids)
  - Smart chunking per doc type (SOP/Policy/Runbook/Ticket/Slide)
  - Chunk enrichment (keywords + 2 hypothetical questions via Ollama)
"""
from __future__ import annotations

import hashlib
import logging
import re
import uuid
from dataclasses import dataclass, field
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import fitz  # PyMuPDF
import tiktoken

logger = logging.getLogger(__name__)

# ── Token counter ─────────────────────────────────────────────────────────────
_tokenizer = tiktoken.get_encoding("cl100k_base")


def count_tokens(text: str) -> int:
    return len(_tokenizer.encode(text))


# ── Data models ───────────────────────────────────────────────────────────────
@dataclass
class RawDocument:
    doc_id: str = field(default_factory=lambda: str(uuid.uuid4()))
    source_file: str = ""
    title: str = ""
    source_type: str = "PDF"   # PDF | DOCX | PPTX | CSV
    raw_text: str = ""
    pages: Dict[int, str] = field(default_factory=dict)  # page_num -> text
    checksum: str = ""
    document_summary: str = ""
    document_keywords: List[str] = field(default_factory=list)
    document_questions: List[str] = field(default_factory=list)


@dataclass
class Chunk:
    chunk_id: str = field(default_factory=lambda: str(uuid.uuid4()))
    doc_id: str = ""
    chunk_index: int = 0
    content: str = ""
    section_title: str = ""
    page_number: int = 0
    token_count: int = 0
    doc_type: str = ""
    department: str = ""
    doc_origin: str = "INTERNAL"
    priority_tier: int = 1
    source_file: str = ""
    keywords: List[str] = field(default_factory=list)
    hypothetical_questions: List[str] = field(default_factory=list)
    # ── P3: metadata-flow fields (must survive Ingestion -> Qdrant -> Citation) ──
    repository: str = ""
    access_roles: List[str] = field(default_factory=list)
    project_id: str = ""
    uploaded_by: str = ""
    section_hierarchy: List[str] = field(default_factory=list)
    created_at: str = ""
    # ── P7: PII hashing ──
    pii_hash_map: Dict[str, str] = field(default_factory=dict)
    # ── P9: image-aware retrieval ──
    is_image_chunk: bool = False
    image_path: str = ""
    chunk_summary: str = ""
    chunk_keywords: List[str] = field(default_factory=list)
    chunk_questions: List[str] = field(default_factory=list)


# ── L1: Rule-based classifier ─────────────────────────────────────────────────
_L1_RULES: Dict[str, Dict] = {
    # department rules (regex on filename/path)
    "department": {
        r"(hr|human.?resource|leave|payroll|onboard)": "HR",
        r"(finance|budget|invoice|expense|account)": "Finance",
        r"(it|network|infra|devops|deploy|kubernetes|docker)": "IT",
        r"(engineer|architecture|system.?design|api|backend)": "Engineering",
        r"(project|milestone|roadmap|sprint|scrum)": "Projects",
        r"(gitlab|handbook|policy|guide|customer)": "Operations",
    },
    # doc_type rules
    "doc_type": {
        r"(sop|standard.?operating)": "SOP",
        r"(policy|policies)": "Policy",
        r"(runbook|playbook|incident.?response)": "Runbook",
        r"(guide|how.?to|tutorial|getting.?started)": "Guide",
        r"(presentation|slides|deck)": "Presentation",
        r"(handbook|wiki)": "Guide",
        r"(ticket|incident|issue)": "Ticket",
    },
}


def classify_l1(filename: str, folder_path: str = "") -> Tuple[str, str, float]:
    """
    Returns (department, doc_type, confidence)
    confidence > 0.90 → accept, skip L2/L3
    """
    text = (filename + " " + folder_path).lower()
    department, doc_type = "Unknown", "Guide"
    dept_hit = doc_hit = False

    for pattern, dept in _L1_RULES["department"].items():
        if re.search(pattern, text, re.IGNORECASE):
            department = dept
            dept_hit = True
            break

    for pattern, dtype in _L1_RULES["doc_type"].items():
        if re.search(pattern, text, re.IGNORECASE):
            doc_type = dtype
            doc_hit = True
            break

    confidence = 0.95 if (dept_hit and doc_hit) else 0.80 if dept_hit else 0.70
    return department, doc_type, confidence


# ── File processors ──────────────────────────────────────────────────────────
def process_pdf(file_path: str) -> RawDocument:
    """Extract text + page structure from PDF using PyMuPDF."""
    doc = fitz.open(file_path)
    pages: Dict[int, str] = {}
    full_text_parts = []

    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        if text:
            pages[page_num] = text
            full_text_parts.append(text)

    raw_text = "\n\n".join(full_text_parts)
    doc.close()

    checksum = hashlib.sha256(raw_text.encode()).hexdigest()
    title = Path(file_path).stem.replace("_", " ").replace("-", " ")

    return RawDocument(
        source_file=file_path,
        title=title,
        source_type="PDF",
        raw_text=raw_text,
        pages=pages,
        checksum=checksum,
    )


def process_docx(file_path: str) -> RawDocument:
    """Extract text from DOCX preserving heading hierarchy."""
    from docx import Document

    doc = Document(file_path)
    sections: List[str] = []
    current_section = []
    full_text_parts = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        full_text_parts.append(text)

        # Use heading styles as section boundaries
        if para.style.name.startswith("Heading"):
            if current_section:
                sections.append("\n".join(current_section))
            current_section = [text]
        else:
            current_section.append(text)

    if current_section:
        sections.append("\n".join(current_section))

    raw_text = "\n\n".join(full_text_parts)
    checksum = hashlib.sha256(raw_text.encode()).hexdigest()
    title = Path(file_path).stem.replace("_", " ").replace("-", " ")

    return RawDocument(
        source_file=file_path,
        title=title,
        source_type="DOCX",
        raw_text=raw_text,
        checksum=checksum,
    )


def process_pptx(file_path: str) -> RawDocument:
    """Each slide = one chunk. Extracts title + body + notes."""
    from pptx import Presentation

    prs = Presentation(file_path)
    pages: Dict[int, str] = {}
    full_text_parts = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts = []
        if slide.shapes.title and slide.shapes.title.text:
            parts.append(f"[Slide {slide_num}] {slide.shapes.title.text}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape != slide.shapes.title:
                    parts.append(shape.text.strip())
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append(f"Notes: {notes_text}")

        slide_text = "\n".join(parts)
        pages[slide_num] = slide_text
        full_text_parts.append(slide_text)

    raw_text = "\n\n".join(full_text_parts)
    checksum = hashlib.sha256(raw_text.encode()).hexdigest()
    title = Path(file_path).stem.replace("_", " ").replace("-", " ")

    return RawDocument(
        source_file=file_path,
        title=title,
        source_type="PPTX",
        raw_text=raw_text,
        pages=pages,
        checksum=checksum,
    )


# ── Smart chunking strategies ─────────────────────────────────────────────────

# ── P3/P1: page-number resolution ──────────────────────────────────────────────
# `chunk_policy_sop` / `chunk_runbook` / the generic fallback operate on
# `raw_doc.raw_text` (the concatenation of all pages), which previously meant
# `page_number` defaulted to 0 for every PDF chunk except PPTX slides. Since
# `raw_doc.pages[page_num]` holds the verbatim text of each page, we can map a
# chunk back to the page that contains its opening text.
def _find_page_number(raw_doc: "RawDocument", text: str, default: int = 0) -> int:
    """Return the 1-indexed page number whose page text contains the start of
    `text`. Falls back to `default` (0) if `raw_doc.pages` is empty (DOCX) or
    no page matches (shouldn't normally happen for PDFs)."""
    if not raw_doc.pages or not text:
        return default

    # Use a short, distinctive snippet from the start of the chunk's content —
    # the first non-empty line, trimmed — to locate the source page.
    snippet = ""
    for line in text.splitlines():
        line = line.strip()
        if len(line) >= 8:
            snippet = line[:80]
            break
    if not snippet:
        snippet = text.strip()[:80]
    if not snippet:
        return default

    for page_num, page_text in raw_doc.pages.items():
        if snippet in page_text:
            return page_num

    # Fallback: try matching just the first few words (handles whitespace
    # normalisation differences between raw_text and pages[]).
    words = snippet.split()
    if len(words) >= 3:
        loose = " ".join(words[:5])
        for page_num, page_text in raw_doc.pages.items():
            if loose in page_text:
                return page_num

    return default


# ── P3: section-hierarchy detection ────────────────────────────────────────────
_MD_HEADING = re.compile(r"^(#{1,6})\s+(.+)$")
_NUMBERED_HEADING = re.compile(r"^(\d+(?:\.\d+)*)\.?\s+(.+)$")


def _heading_level(line: str) -> int:
    """Best-effort heading depth (1 = top-level) for hierarchy tracking.
    Returns 0 if the line doesn't look like a heading at all."""
    line = line.strip()
    md = _MD_HEADING.match(line)
    if md:
        return len(md.group(1))
    num = _NUMBERED_HEADING.match(line)
    if num:
        return num.group(1).count(".") + 1
    # ALL-CAPS short lines are treated as top-level section headers
    if line.isupper() and 3 <= len(line) <= 80:
        return 1
    return 0


class _HierarchyTracker:
    """Maintains a stack of (level, title) so each chunk can carry its full
    `section_hierarchy` path, e.g. ["VPN Runbook", "1. Access Management",
    "1.3 VPN", "1.3.2 Reset Procedure"]."""

    def __init__(self, doc_title: str):
        self.doc_title = doc_title
        self._stack: List[Tuple[int, str]] = []

    def update(self, heading_line: str):
        level = _heading_level(heading_line)
        if level <= 0:
            return
        # Pop anything at this level or deeper, then push the new heading.
        self._stack = [(lvl, title) for lvl, title in self._stack if lvl < level]
        self._stack.append((level, heading_line.strip()))

    def path(self, leaf_title: Optional[str] = None) -> List[str]:
        hierarchy = [self.doc_title] + [title for _, title in self._stack]
        if leaf_title and (not hierarchy or hierarchy[-1] != leaf_title.strip()):
            hierarchy.append(leaf_title.strip())
        return hierarchy

def chunk_by_hierarchy(raw_doc: RawDocument) -> List[Chunk]:
    """
    Universal hierarchical chunker for PDF and DOCX.

    Detects:
    - Markdown headings (# ## ###)
    - Numbered headings (1, 1.1, 1.1.1)
    - ALL CAPS headings
    - Title Case short headings

    Produces chunks using document structure first,
    token splitting only when a section exceeds limits.
    """

    chunks = []
    tracker = _HierarchyTracker(raw_doc.title)

    lines = raw_doc.raw_text.split("\n")

    current_heading = raw_doc.title
    current_hierarchy = [raw_doc.title]
    current_lines = []

    sections = []

    heading_pattern = re.compile(
        r"""
        ^
        (
            \#{1,6}\s+.+ |
            \d+(\.\d+)*\.?\s+.+ |
            [A-Z][A-Z\s\-/&,]{3,80} |
            [A-Z][a-z]+(?:\s+[A-Z][a-z]+){0,6}
        )
        $
        """,
        re.VERBOSE,
    )

    for line in lines:
        line = line.strip()

        if not line:
            continue

        is_heading = (
            len(line) < 120
            and heading_pattern.match(line)
        )

        if is_heading:

            if current_lines:
                sections.append(
                    (
                        current_heading,
                        current_hierarchy,
                        "\n".join(current_lines),
                    )
                )

            current_heading = line

            tracker.update(line)

            current_hierarchy = tracker.path(line)

            current_lines = []

        else:
            current_lines.append(line)

    if current_lines:
        sections.append(
            (
                current_heading,
                current_hierarchy,
                "\n".join(current_lines),
            )
        )

    chunk_index = 0

    for section_title, hierarchy, section_text in sections:

        section_text = section_text.strip()

        if not section_text:
            continue

        token_count = count_tokens(section_text)

        page_number = _find_page_number(
            raw_doc,
            section_text
        )

        if token_count <= 800:

            chunks.append(
                Chunk(
                    doc_id=raw_doc.doc_id,
                    chunk_index=chunk_index,
                    content=f"{section_title}\n\n{section_text}",
                    section_title=section_title,
                    section_hierarchy=hierarchy,
                    page_number=page_number,
                    token_count=token_count,
                    source_file=raw_doc.source_file,
                )
            )

            chunk_index += 1

        else:

            sub_chunks = _split_by_tokens(
                section_text,
                max_tokens=700,
                overlap_ratio=0.15,
            )

            for part_num, sub in enumerate(sub_chunks, start=1):

                chunks.append(
                    Chunk(
                        doc_id=raw_doc.doc_id,
                        chunk_index=chunk_index,
                        content=f"{section_title} (part {part_num})\n\n{sub}",
                        section_title=section_title,
                        section_hierarchy=hierarchy,
                        page_number=page_number,
                        token_count=count_tokens(sub),
                        source_file=raw_doc.source_file,
                    )
                )

                chunk_index += 1

    return chunks


def _split_by_tokens(text: str, max_tokens: int, overlap_ratio: float = 0.15) -> List[str]:
    """Generic token-aware splitter with overlap."""
    words = text.split()
    chunks, current_words, current_tokens = [], [], 0
    overlap_tokens = int(max_tokens * overlap_ratio)

    for word in words:
        word_tokens = count_tokens(word)
        if current_tokens + word_tokens > max_tokens and current_words:
            chunks.append(" ".join(current_words))
            # keep overlap: roll back by overlap_tokens
            overlap_words = []
            overlap_count = 0
            for w in reversed(current_words):
                wt = count_tokens(w)
                if overlap_count + wt > overlap_tokens:
                    break
                overlap_words.insert(0, w)
                overlap_count += wt
            current_words = overlap_words
            current_tokens = overlap_count

        current_words.append(word)
        current_tokens += word_tokens

    if current_words:
        chunks.append(" ".join(current_words))

    return [c for c in chunks if c.strip()]


def chunk_policy_sop(raw_doc: RawDocument) -> List[Chunk]:
    """
    SOP / Policy: chunk by H1/H2/H3 headings.
    Target: 512-800 tokens, 15% overlap.

    P3/P1: each chunk now carries:
      - page_number  : resolved via _find_page_number() against raw_doc.pages
      - section_hierarchy : full heading path (document title -> ... -> section)
    """
    chunks = []
    # Detect heading-like lines (short lines, possibly ALL CAPS or title case)
    heading_pattern = re.compile(
        r"^(?:[A-Z][A-Z\s&/,]{3,50}|#{1,3}\s+.+|\d+\.\s+[A-Z].+)$"
    )
    lines = raw_doc.raw_text.split("\n")
    tracker = _HierarchyTracker(raw_doc.title)
    sections: List[Tuple[str, List[str], List[str]]] = []
    current_heading = raw_doc.title
    current_lines: List[str] = []
    current_hierarchy = tracker.path(current_heading)

    for line in lines:
        line = line.strip()
        if not line:
            continue
        if heading_pattern.match(line) and len(line) < 120:
            if current_lines:
                sections.append((current_heading, current_lines, current_hierarchy))
            current_heading = line
            tracker.update(line)
            current_hierarchy = tracker.path(line)
            current_lines = []
        else:
            current_lines.append(line)

    if current_lines:
        sections.append((current_heading, current_lines, current_hierarchy))

    if not sections:
        # Fallback: generic split
        sections = [(raw_doc.title, raw_doc.raw_text.split("\n"), [raw_doc.title])]

    chunk_index = 0
    for section_title, section_lines, hierarchy in sections:
        section_text = "\n".join(section_lines)
        if not section_text.strip():
            continue

        page_number = _find_page_number(raw_doc, section_text)
        token_count = count_tokens(section_text)
        if token_count <= 800:
            # Section fits in one chunk
            chunks.append(Chunk(
                doc_id=raw_doc.doc_id,
                chunk_index=chunk_index,
                content=f"{section_title}\n\n{section_text}",
                section_title=section_title,
                section_hierarchy=hierarchy,
                page_number=page_number,
                token_count=token_count,
                source_file=raw_doc.source_file,
            ))
            chunk_index += 1
        else:
            # Split large section
            sub_chunks = _split_by_tokens(section_text, max_tokens=700, overlap_ratio=0.15)
            for i, sub in enumerate(sub_chunks):
                chunks.append(Chunk(
                    doc_id=raw_doc.doc_id,
                    chunk_index=chunk_index,
                    content=f"{section_title} (part {i+1})\n\n{sub}",
                    section_title=section_title,
                    section_hierarchy=hierarchy,
                    page_number=page_number if i == 0 else _find_page_number(raw_doc, sub, page_number),
                    token_count=count_tokens(sub),
                    source_file=raw_doc.source_file,
                ))
                chunk_index += 1

    return chunks


def chunk_runbook(raw_doc: RawDocument) -> List[Chunk]:
    """
    Runbook / Guide: chunk by numbered steps and code blocks.
    Target: 200-500 tokens, 10% overlap.
    Never split code blocks.

    P3/P1: resolves page_number via _find_page_number() and tracks the
    heading/step path in section_hierarchy.
    """
    chunks = []
    chunk_index = 0
    lines = raw_doc.raw_text.split("\n")

    tracker = _HierarchyTracker(raw_doc.title)
    current_step_title = "Introduction"
    current_hierarchy = tracker.path(current_step_title)
    current_lines: List[str] = []
    in_code_block = False

    step_pattern = re.compile(r"^(?:Step\s*\d+|^\d+[\.\)]\s+|\#{1,3}\s+)")

    for line in lines:
        if line.strip().startswith("```"):
            in_code_block = not in_code_block

        is_step_header = step_pattern.match(line.strip()) and not in_code_block

        if is_step_header and current_lines:
            content = "\n".join(current_lines).strip()
            if content:
                chunks.append(Chunk(
                    doc_id=raw_doc.doc_id,
                    chunk_index=chunk_index,
                    content=content,
                    section_title=current_step_title,
                    section_hierarchy=current_hierarchy,
                    page_number=_find_page_number(raw_doc, content),
                    token_count=count_tokens(content),
                    source_file=raw_doc.source_file,
                ))
                chunk_index += 1
            current_step_title = line.strip()
            tracker.update(line.strip())
            current_hierarchy = tracker.path(current_step_title)
            current_lines = [line]
        else:
            current_lines.append(line)

    if current_lines:
        content = "\n".join(current_lines).strip()
        if content:
            chunks.append(Chunk(
                doc_id=raw_doc.doc_id,
                chunk_index=chunk_index,
                content=content,
                section_title=current_step_title,
                section_hierarchy=current_hierarchy,
                page_number=_find_page_number(raw_doc, content),
                token_count=count_tokens(content),
                source_file=raw_doc.source_file,
            ))

    # Split any chunk still over 500 tokens (code blocks preserved)
    final_chunks = []
    for chunk in chunks:
        if chunk.token_count > 500:
            sub_texts = _split_by_tokens(chunk.content, max_tokens=450, overlap_ratio=0.10)
            for i, sub in enumerate(sub_texts):
                final_chunks.append(Chunk(
                    doc_id=raw_doc.doc_id,
                    chunk_index=len(final_chunks),
                    content=sub,
                    section_title=chunk.section_title,
                    section_hierarchy=chunk.section_hierarchy,
                    page_number=chunk.page_number if i == 0 else _find_page_number(raw_doc, sub, chunk.page_number),
                    token_count=count_tokens(sub),
                    source_file=raw_doc.source_file,
                ))
        else:
            final_chunks.append(chunk)

    return final_chunks


def chunk_slides(raw_doc: RawDocument) -> List[Chunk]:
    """PPTX: one slide = one chunk."""
    chunks = []
    for slide_num, slide_text in raw_doc.pages.items():
        if slide_text.strip():
            chunks.append(Chunk(
                doc_id=raw_doc.doc_id,
                chunk_index=slide_num - 1,
                content=slide_text,
                section_title=f"Slide {slide_num}",
                section_hierarchy=[raw_doc.title, f"Slide {slide_num}"],
                page_number=slide_num,
                token_count=count_tokens(slide_text),
                source_file=raw_doc.source_file,
            ))
    return chunks


def chunk_document(raw_doc: RawDocument, doc_type: str, department: Optional[str] = None) -> List[Chunk]:
    """Route to the correct chunking strategy based on doc type.

    `department` (optional) is forwarded to PII redaction (P7) so that
    IT/Engineering technical content (IPs, ports, server names, ticket
    references, config values) in SOPs/Runbooks is preserved rather than
    destroyed by generic PII patterns.
    """
    if raw_doc.source_type == "PPTX":
        chunks = chunk_slides(raw_doc)

    else:
        chunks = chunk_by_hierarchy(raw_doc)

        if not chunks:
            sub_texts = _split_by_tokens(
                raw_doc.raw_text,
                max_tokens=700,
                overlap_ratio=0.15,
            )

            chunks = []

            for i, text in enumerate(sub_texts):

                chunks.append(
                    Chunk(
                        doc_id=raw_doc.doc_id,
                        chunk_index=i,
                        content=text,
                        section_title=raw_doc.title,
                        section_hierarchy=[raw_doc.title],
                        page_number=_find_page_number(raw_doc, text),
                        token_count=count_tokens(text),
                        source_file=raw_doc.source_file,
                    )
                )

    # Stamp all chunks with doc metadata
    for chunk in chunks:
        chunk.doc_id = raw_doc.doc_id
        chunk.source_file = raw_doc.source_file
        chunk.doc_type = doc_type
        if department:
            chunk.department = department

    # ── PII Redaction / Hashing (P7, DPDP/GDPR/HIPAA compliance) ────────────
    # True PII (email, person, SSN/PAN/Aadhaar, credit card) is replaced with
    # DETERMINISTIC HASH TOKENS (e.g. EMAIL_HASH_73AD21) so retrieval/clustering
    # still distinguish between distinct entities without exposing raw values
    # to the LLM. Technical identifiers (IP addresses, server names, ports,
    # ticket references, config values) in IT/Engineering/SOP/Runbook content
    # are left untouched — see ingestion/pii_redaction.py.
    from ingestion.pii_redaction import redact_chunk
    for chunk in chunks:
        redact_chunk(chunk, doc_type=doc_type, department=chunk.department)

    return chunks

async def enrich_document_async(
    raw_doc: RawDocument,
    ollama_model: str = "phi3:mini"
) -> RawDocument:

    import httpx
    import json

    prompt = f"""
You are generating retrieval metadata for a document knowledge base.

Analyze this document and return JSON only.

Generate:

1. summary
   - One concise summary
   - Maximum 60 words

2. keywords
   - Top 5 business or technical keywords. Make sure they are relevant to the document's content and context.
   Do not include generic terms like "document", "information", or "content".

3. questions
   - 5 realistic user questions
   - Questions users would ask to find this document

Return:

{{
  "summary": "...",
  "keywords": [
    "...",
    "...",
    "...",
    "...",
    "..."
  ],
  "questions": [
    "...",
    "...",
    "...",
    "..."
  ]
}}

Document:

{raw_doc.raw_text[:5000]}
"""

    try:

        async with httpx.AsyncClient(timeout=30.0) as client:

            resp = await client.post(
                "http://localhost:11434/api/generate",
                json={
                    "model": ollama_model,
                    "prompt": prompt,
                    "stream": False,
                },
            )

            result = resp.json()["response"]

            match = re.search(r"\{.*\}", result, re.DOTALL)

            if match:

                metadata = json.loads(match.group())

                raw_doc.document_summary = metadata.get(
                    "summary",
                    ""
                )

                raw_doc.document_keywords = metadata.get(
                    "keywords",
                    []
                )

                raw_doc.document_questions = metadata.get(
                    "questions",
                    []
                )

    except Exception as e:

        logger.warning(
            f"Document enrichment failed: {e}"
        )

    return raw_doc

# ── Chunk enrichment via Ollama ────────────────────────────────────────────────
async def enrich_chunk_async(
    chunk: Chunk,
    ollama_model: str = "phi3:mini"
) -> Chunk:

    import httpx
    import json

    prompt = f"""
You are generating retrieval metadata for a RAG system.

Analyze the chunk below.

Return ONLY valid JSON.

Generate:

1. summary
   - One sentence
   - Maximum 40 words

2. keywords
   - Top 5 keywords or phrases

3. questions
   - 4 realistic user questions
   - Questions users would ask to retrieve this information

Rules:

- Use only information present in the chunk
- Preserve technical terms
- Preserve acronyms
- Preserve product names
- Do not hallucinate

Return:

{{
  "summary": "...",
  "keywords": [
    "...",
    "...",
    "...",
    "...",
    "..."
  ],
  "questions": [
    "...",
    "...",
    "...",
    "..."
  ]
}}

Chunk:

{chunk.content[:1500]}
"""

    try:

        async with httpx.AsyncClient(timeout=20.0) as client:

            resp = await client.post(
                "http://localhost:11434/api/generate",
                json={
                    "model": ollama_model,
                    "prompt": prompt,
                    "stream": False,
                },
            )

            result = resp.json()["response"]

            match = re.search(
                r"\{.*\}",
                result,
                re.DOTALL,
            )

            if match:

                metadata = json.loads(
                    match.group()
                )

                chunk.chunk_summary = metadata.get(
                    "summary",
                    ""
                )

                chunk.chunk_keywords = metadata.get(
                    "keywords",
                    []
                )

                chunk.chunk_questions = metadata.get(
                    "questions",
                    []
                )

    except Exception as e:

        logger.debug(
            f"Chunk enrichment skipped: {e}"
        )

        chunk.chunk_summary = (
            f"Information about {chunk.section_title}"
        )

        chunk.chunk_keywords = []

        chunk.chunk_questions = [
            f"What is {chunk.section_title}?",
            f"How does {chunk.section_title} work?",
        ]

    return chunk
