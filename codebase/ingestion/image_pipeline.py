"""
ingestion/image_pipeline.py — P9: Image-Aware Retrieval

Pipeline:  Image -> OCR -> Caption -> Embedding -> Searchable Chunk

Extracts embedded images (PDF page images, PPTX picture shapes, DOCX inline
images) and converts each into a `Chunk` whose `content` = caption + OCR
text. Because these chunks flow through the SAME `vector_store.upsert_chunks`
/ `bm25_store.add_documents` path as text chunks, a query like:

    "Show me the network architecture diagram"

retrieves them via the existing content/question vector search — no separate
image index or UI changes required.

Graceful degradation (CPU-only, network-optional):
  - OCR (pytesseract) is optional. If unavailable, ocr_image() returns "".
  - Captioning uses a vision-capable Ollama model (qwen2.5-vl). If Ollama or
    the model is unavailable, falls back to an OCR-text-based or
    context-based caption so ingestion never hard-fails.
  - Tiny embedded images (icons, bullets, logos) below `min_image_bytes` are
    skipped so they don't pollute retrieval with noise chunks.
"""
from __future__ import annotations

import base64
import logging
import os
from dataclasses import dataclass, field
from typing import List, Optional

import httpx

from configs.settings import get_settings
from ingestion.idp_pipeline import Chunk, RawDocument, count_tokens

logger = logging.getLogger(__name__)
settings = get_settings()


@dataclass
class ExtractedImage:
    image_bytes: bytes
    ext: str = "png"
    page_number: int = 0
    # Nearby text (slide title, surrounding paragraph, page text) used to
    # give the vision captioner context about what the image likely shows.
    context_hint: str = ""


# ════════════════════════════════════════════════════════════════════════════
# EXTRACTION — PDF / PPTX / DOCX
# ════════════════════════════════════════════════════════════════════════════
def extract_images_from_pdf(file_path: str) -> List[ExtractedImage]:
    """Extract embedded images from a PDF using PyMuPDF (covers PDF page
    images, embedded architecture diagrams, screenshots, etc.)."""
    import fitz

    images: List[ExtractedImage] = []
    doc = fitz.open(file_path)
    for page_num, page in enumerate(doc, start=1):
        page_text = page.get_text("text")[:300].strip()
        for img in page.get_images(full=True):
            xref = img[0]
            try:
                base_image = doc.extract_image(xref)
                images.append(ExtractedImage(
                    image_bytes=base_image["image"],
                    ext=base_image.get("ext", "png"),
                    page_number=page_num,
                    context_hint=page_text,
                ))
            except Exception as e:
                logger.debug(f"PDF image extract failed (page={page_num}, xref={xref}): {e}")
    return images


def extract_images_from_pptx(file_path: str) -> List[ExtractedImage]:
    """Extract picture shapes from a PPTX (covers screenshots, architecture
    diagrams pasted into slides)."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    images: List[ExtractedImage] = []
    prs = Presentation(file_path)
    for slide_num, slide in enumerate(prs.slides, start=1):
        title = ""
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()

        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                image = shape.image
                images.append(ExtractedImage(
                    image_bytes=image.blob,
                    ext=image.ext or "png",
                    page_number=slide_num,
                    context_hint=title or f"Slide {slide_num}",
                ))
            except Exception as e:
                logger.debug(f"PPTX image extract failed (slide={slide_num}): {e}")
    return images


def extract_images_from_docx(file_path: str) -> List[ExtractedImage]:
    """Extract inline images embedded in a DOCX via its package relationships."""
    from docx import Document

    images: List[ExtractedImage] = []
    doc = Document(file_path)
    title_hint = doc.paragraphs[0].text.strip() if doc.paragraphs else ""

    for rel in doc.part.rels.values():
        if "image" not in rel.reltype:
            continue
        try:
            blob = rel.target_part.blob
            content_type = getattr(rel.target_part, "content_type", "image/png")
            ext = content_type.split("/")[-1] or "png"
            images.append(ExtractedImage(
                image_bytes=blob,
                ext=ext,
                page_number=0,  # DOCX has no native page concept (python-docx)
                context_hint=title_hint,
            ))
        except Exception as e:
            logger.debug(f"DOCX image extract failed: {e}")
    return images


_EXTRACTORS = {
    ".pdf": extract_images_from_pdf,
    ".pptx": extract_images_from_pptx,
    ".docx": extract_images_from_docx,
}


# ════════════════════════════════════════════════════════════════════════════
# OCR
# ════════════════════════════════════════════════════════════════════════════
def ocr_image(image_bytes: bytes) -> str:
    """Extract any text rendered inside the image (diagram labels, annotated
    screenshots). Returns "" if pytesseract/Pillow are unavailable or OCR
    finds nothing — never raises."""
    try:
        import io
        from PIL import Image
        import pytesseract

        img = Image.open(io.BytesIO(image_bytes))
        text = pytesseract.image_to_string(img)
        return text.strip()
    except Exception as e:
        logger.debug(f"OCR unavailable/failed: {e}")
        return ""


# ════════════════════════════════════════════════════════════════════════════
# CAPTIONING (vision model via Ollama)
# ════════════════════════════════════════════════════════════════════════════
_CAPTION_PROMPT = (
    "Describe this image in 1-3 sentences for a document search index. "
    "If it is a diagram, architecture diagram, flowchart, screenshot, or "
    "chart, say so explicitly and describe what it shows (components, "
    "labels, relationships). Be factual and specific."
)


async def caption_image_async(
    image_bytes: bytes, ocr_text: str = "", context_hint: str = ""
) -> str:
    """
    Caption an image using a vision-capable Ollama model (qwen2.5-vl).

    Falls back to an OCR-text-based or context-based caption if the vision
    model / Ollama is unavailable — captioning is a quality enhancement, not
    a hard ingestion dependency.
    """
    try:
        b64 = base64.b64encode(image_bytes).decode("ascii")
        prompt = _CAPTION_PROMPT
        if ocr_text:
            prompt += f"\n\nText detected in the image (OCR): {ocr_text[:300]}"
        if context_hint:
            prompt += f"\n\nSurrounding document context: {context_hint[:200]}"

        async with httpx.AsyncClient(timeout=20.0) as client:
            resp = await client.post(
                f"{settings.ollama_base_url}/api/generate",
                json={
                    "model": settings.vision_model,
                    "prompt": prompt,
                    "images": [b64],
                    "stream": False,
                    "options": {"temperature": 0.2, "num_predict": 150},
                },
            )
            data = resp.json()
            caption = data.get("response", "").strip()
            if caption:
                return caption
    except Exception as e:
        logger.debug(f"Vision captioning unavailable ({settings.vision_model}): {e}")

    # ── Fallback captions (no vision model available) ──────────────────────
    if ocr_text:
        return f"Image containing text: {ocr_text[:200]}"
    if context_hint:
        return f"Embedded image near: {context_hint[:150]}"
    return "Embedded image (no caption available)"


# ════════════════════════════════════════════════════════════════════════════
# CHUNK CONSTRUCTION
# ════════════════════════════════════════════════════════════════════════════
def build_image_chunk(
    raw_doc: RawDocument,
    extracted: ExtractedImage,
    index: int,
    ocr_text: str,
    caption: str,
    save_dir: str,
) -> Chunk:
    """Persist the image to disk and build a Chunk whose `content` (caption
    + OCR text) is what gets embedded/indexed — making the image
    semantically searchable like any text chunk."""
    os.makedirs(save_dir, exist_ok=True)
    ext = (extracted.ext or "png").lstrip(".")
    filename = f"{raw_doc.doc_id}_img{index}.{ext}"
    path = os.path.join(save_dir, filename)
    try:
        with open(path, "wb") as f:
            f.write(extracted.image_bytes)
    except Exception as e:
        logger.warning(f"Failed to persist extracted image {filename}: {e}")
        path = ""

    content_parts = [caption] if caption else []
    if ocr_text:
        content_parts.append(f"Text extracted from image (OCR): {ocr_text}")
    if not content_parts:
        content_parts.append(f"Embedded image from {raw_doc.title}")
    content = "\n\n".join(content_parts)

    if extracted.page_number:
        section_title = f"Image (page {extracted.page_number})"
    else:
        section_title = "Embedded Image"

    return Chunk(
        doc_id=raw_doc.doc_id,
        content=content,
        section_title=section_title,
        section_hierarchy=[raw_doc.title, section_title],
        page_number=extracted.page_number,
        token_count=count_tokens(content),
        source_file=raw_doc.source_file,
        doc_type="Diagram",
        is_image_chunk=True,
        image_path=path,
    )


# ════════════════════════════════════════════════════════════════════════════
# PUBLIC ENTRY POINT
# ════════════════════════════════════════════════════════════════════════════
async def extract_image_chunks(
    file_path: str, raw_doc: RawDocument, ext: str, save_dir: Optional[str] = None,
) -> List[Chunk]:
    """
    Run the full Image -> OCR -> Caption -> Chunk pipeline for one document.

    `ext` is the file extension (".pdf" / ".pptx" / ".docx"). Returns an
    empty list for unsupported extensions or if extraction fails — image
    extraction failure is never fatal to document ingestion.
    """
    extractor = _EXTRACTORS.get(ext.lower())
    if extractor is None:
        return []

    save_dir = save_dir or os.path.join(settings.image_storage_dir, raw_doc.doc_id)

    try:
        extracted_images = extractor(file_path)
    except Exception as e:
        logger.warning(f"Image extraction failed for {file_path}: {e}")
        return []

    chunks: List[Chunk] = []
    for i, img in enumerate(extracted_images):
        if len(img.image_bytes) < settings.min_image_bytes:
            continue  # skip icons/bullets/logos

        ocr_text = ocr_image(img.image_bytes)
        caption = await caption_image_async(img.image_bytes, ocr_text, img.context_hint)
        chunks.append(build_image_chunk(raw_doc, img, i, ocr_text, caption, save_dir))

    if chunks:
        logger.info(
            "[INGEST DEBUG] Image pipeline: %d image chunk(s) extracted from %s",
            len(chunks), os.path.basename(file_path),
        )
    return chunks
